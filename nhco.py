from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION, XL_TICK_MARK
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd


def load_clean_csv(path):
    df = pd.read_csv(path)
    if "row_number" in df.columns:
        df = df.drop(columns=["row_number"])
    return df

font_standard = "Gotham HTF"
font_light = "Gotham HTF Light" 
font_medium = "Gotham HTF Medium"

files_folder = "/files"
template = f"{files_folder}/NHCO_template.pptx"

font_standard = "Gotham HTF"
font_light = "Gotham HTF Light" 
font_medium = "Gotham HTF Medium"

timing_media_spending_data = load_clean_csv(f"{files_folder}/nhco_timing_media_spending.csv")
dpa_table_data = load_clean_csv(f"{files_folder}/nhco_dpa_tabella.csv")
dpa_impression_grafico_data = load_clean_csv(f"{files_folder}/nhco_dpa_impression_grafico.csv")
dpa_viewability_grafico_data = load_clean_csv(f"{files_folder}/nhco_dpa_viewability_grafico.csv")

# da attivare
# vpa_table_data = load_clean_csv(f"{files_folder}/nhco_vpa_tabella.csv")
# vpa_impression_grafico_data = load_clean_csv(f"{files_folder}/nhco_vpa_impression_grafico.csv")
# vpa_completion_grafico_data = load_clean_csv(f"{files_folder}/nhco_vpa_completion_grafico.csv")


from pptx.oxml import parse_xml
from pptx.oxml.ns import qn

def set_chart_overlap_and_gap(chart, overlap_value, gap_width_value):
    chart_xml = chart._element

    # Trova l'elemento <c:barChart>
    bar_charts = chart_xml.findall(".//" + qn("c:barChart"))
    if not bar_charts:
        raise ValueError("Nessun elemento <c:barChart> trovato nel grafico")
    bar_chart = bar_charts[0]

    # Crea ed inserisce overlap
    overlap_xml = parse_xml(f'<c:overlap xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" val="{overlap_value}"/>')
    bar_chart.insert(0, overlap_xml)

    # Crea ed inserisce gapWidth
    gap_width_xml = parse_xml(f'<c:gapWidth xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" val="{gap_width_value}"/>')
    bar_chart.insert(1, gap_width_xml)


from pptx import Presentation

def replace_text_in_slide(slide, replacements):
    """
    Sostituisce i segnaposto nelle forme di testo della slide.

    :param slide: oggetto Slide
    :param replacements: dict con chiavi=segnaposto, valori=testo sostitutivo
                         es. {"{{TITOLO_REPORT}}": "Report Mensile Maggio"}
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for placeholder, new_text in replacements.items():
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, new_text)


def fill_timing_media_spending_table(slide, data):
    """
    Riempie la tabella presente nella slide con i dati del CSV.
    """
    # Trova la tabella
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break
    
    if table is None:
        raise ValueError("Nessuna tabella trovata nella slide.")
    
    # Inserisci i dati nelle celle della tabella
    if len(table.rows) < len(data) + 1:
        raise ValueError(f"La tabella ha solo {len(table.rows)} righe, ma i dati richiedono {len(data) + 1} righe.")

    if len(table.columns) < len(data.columns):
        raise ValueError(f"La tabella ha solo {len(table.columns)} colonne, ma i dati richiedono {len(data.columns)} colonne.")

    for i, row in data.iterrows():
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)  # +1 perché riga 0 è intestazione
            cell.text = str(value)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)

            # Applica font_standard solo all'ultima riga
            if i == len(data) - 1:
                paragraph.font.name = font_standard
                paragraph.font.bold = True
            else:
                paragraph.font.name = font_light

            # Centra il testo
            paragraph.alignment = PP_ALIGN.CENTER


def fill_programmatic_table(slide, data):
    """
    Riempie la tabella presente nella slide con i dati del CSV.
    """
    # Trova la tabella
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break
    
    if table is None:
        raise ValueError("Nessuna tabella trovata nella slide.")
    
    # Inserisci i dati nelle celle della tabella a partire dalla riga 1 (0 è intestazione) e colonna 1 (0 è intestazione)
    for i, row in data.iterrows():
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j + 1)  # +1 perché riga 0 è intestazione e +1 perché colonna 0 è intestazione
            cell.text = str(value)
            # centra il testo nelle celle
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_dpa_chart(data, slide, chart_type):
    config = {
        "impression": {
            "value_col": "Impression",
            "target_col": "target",
            "series_names": ["Impression", "target"],
            "title": "Impression",
            "color_map": {
                "Impression": RGBColor(0, 153, 76),
                "target": RGBColor(197, 224, 180),
            },
            "y_min": 0,
            "y_max": 1000000,
            "position": (Cm(1.63), Cm(12.5), Cm(13.88), Cm(4.43)),
        },
        "impression_video": {
            "value_col": "Impression",
            "target_col": "target",
            "series_names": ["Impression", "target"],
            "title": "Impression",
            "color_map": {
                "Impression": RGBColor(0, 153, 76),
                "target": RGBColor(197, 224, 180),
            },
            "y_min": 0,
            "y_max": 200000,
            "position": (Cm(1.63), Cm(12.77), Cm(13.88), Cm(4.43)),
        },
        "viewability": {
            "value_col": "Viewability rate",
            "target_col": "target",
            "series_names": ["Viewability rate", "target"],
            "title": "Viewability rate",
            "color_map": {
                "Viewability rate": RGBColor(0, 153, 76),
                "target": RGBColor(197, 224, 180),
            },
            "y_min": 0.0,
            "y_max": 1.0,
            "position": (Cm(16.25), Cm(12.69), Cm(13.88), Cm(4.43)),
        },
        "completion": {
            "value_col": "Completion rate",
            "target_col": "target",
            "series_names": ["Completion rate", "target"],
            "title": "Completion rate",
            "color_map": {
                "Completion rate": RGBColor(0, 153, 76),
                "target": RGBColor(197, 224, 180),
            },
            "y_min": 0.6,
            "y_max": 0.9,
            "position": (Cm(16.25), Cm(12.77), Cm(13.88), Cm(4.43)),
        }
    }

    conf = config[chart_type]

    # Preprocessing
    df = data.rename(columns={"Month": "mese"})
    df[conf["value_col"]] = df[conf["value_col"]].fillna(0)
    df[conf["target_col"]] = df[conf["target_col"]].fillna(0)

    # Prepare chart data
    chart_data = CategoryChartData()
    chart_data.categories = df["mese"]
    chart_data.add_series(conf["series_names"][0], df[conf["value_col"]])
    chart_data.add_series(conf["series_names"][1], df[conf["target_col"]])

    # Add chart
    x, y, cx, cy = conf["position"]
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart
    set_chart_overlap_and_gap(chart, overlap_value=-29, gap_width_value=219)
    chart.font.name = font_light

    # Set colors
    for series in chart.series:
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = conf["color_map"].get(series.name, RGBColor(0, 0, 0))

    # Y Axis
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.maximum_scale = conf["y_max"]
    chart.value_axis.minimum_scale = conf["y_min"]
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.font_name = font_light
    chart.value_axis.format.line.fill.background()

    gridlines = chart.value_axis.major_gridlines
    line = gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(200, 200, 200)

    # X Axis
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(10)
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.minor_tick_mark = XL_TICK_MARK.NONE
    category_axis.format.line.fill.solid()
    category_axis.format.line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    category_axis.tick_labels.font.font_name = font_light

    # Legend
    chart.has_legend = True
    chart.legend.font.size = Pt(10)
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.font_name = font_light

    # Title
    chart.has_title = True
    chart.chart_title.text_frame.text = conf["title"]
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.name = font_medium


prs = Presentation(template)
slides = prs.slides

replacements = {
    "{{TITOLO_REPORT}}": "REPORT MARZO 2025",
    "{{SITO1_POSIZIONAMENTO_DISPLAY}}": "ANSA.IT",
    "{{SITO2_POSIZIONAMENTO_DISPLAY}}": "LANAZIONE.IT",
    "{{SITO1_POSIZIONAMENTO_VIDEO}}": "ILMETEO.IT",
    "{{SITO2_POSIZIONAMENTO_VIDEO}}": "ILMESSAGGERO.IT",
}

for slide in slides:
    replace_text_in_slide(slide, replacements)

fill_timing_media_spending_table(slides[2], timing_media_spending_data)
fill_programmatic_table(slides[3], dpa_table_data)
add_dpa_chart(dpa_impression_grafico_data, slides[3], chart_type="impression")
add_dpa_chart(dpa_viewability_grafico_data, slides[3], chart_type="viewability")
#fill_programmatic_table(slides[5], vpa_table_data)
#add_dpa_chart(vpa_impression_grafico_data, slides[5], chart_type="impression_video")
#add_dpa_chart(vpa_completion_grafico_data, slides[5], chart_type="completion")

prs.save(f"{files_folder}/nhco_output_presentazione_con_dati.pptx")

