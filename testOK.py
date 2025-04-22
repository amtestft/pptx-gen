import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- Caricamento dati ---
df_sessions = pd.read_csv("/data/sessions.csv")
df_brand = pd.read_csv("/data/top10_brand_micro_moments.csv")
df_product = pd.read_csv("/data/top10_product_micro_moments.csv")
df_need = pd.read_csv("/data/top10_need_micro_moments.csv")
df_best = pd.read_csv("/data/top10_best_volume_micro_moments.csv")

# --- Preprocessing ---
def classify_channel(channel):
    if pd.isna(channel):
        return 'Other'
    if 'Organic' in channel:
        return 'Organic'
    elif 'Paid' in channel:
        return 'Paid'
    else:
        return 'Other'

df_sessions['channel_group'] = df_sessions['Channel'].apply(classify_channel)
df_sessions['month'] = df_sessions['Month'].astype(str)
grouped = df_sessions.groupby(['month', 'channel_group'])['Sessions'].sum().unstack(fill_value=0)
grouped = grouped.sort_index()

COLUMN_RENAMES = {
    "sum_sessions": "SESSIONS",
    "month": "MESE",
    "sessionDefaultChannelGroup": "CHANNEL",
    "impressions": "IMPRESSIONS",
    "micro_moment": "MICRO-MOMENT"
}
COLUMNS_TO_DROP = ["channel_group", "type"]
METRIC_COLUMNS = ["SESSIONS", "IMPRESSIONS"]

def prepare_table(df):
    df = df.copy()
    df.rename(columns=COLUMN_RENAMES, inplace=True)
    df.drop(columns=[col for col in COLUMNS_TO_DROP if col in df.columns], inplace=True)
    for col in df.select_dtypes(include='number').columns:
        df[col] = df[col].map(lambda x: round(x, 1))
    cols = [col for col in df.columns if col not in METRIC_COLUMNS]
    cols += [col for col in METRIC_COLUMNS if col in df.columns]
    return df[cols]

def add_title_slide(prs, title, subtitle=None, black_bg=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if black_bg:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = textbox.text_frame
    tf.text = title
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) if black_bg else RGBColor(0, 0, 0)

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255) if black_bg else RGBColor(0, 0, 0)

def add_blank_slide(prs, black_bg=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if black_bg:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)
    return slide

def add_formatted_table(slide, df):
    table_rows = len(df) + 1
    table_cols = len(df.columns)
    col_width = 9 / table_cols
    row_height = 5 / table_rows
    left = Inches((10 - (col_width * table_cols)) / 2)
    top = Inches(1)

    table_shape = slide.shapes.add_table(
        table_rows, table_cols, left, top,
        Inches(col_width * table_cols), Inches(row_height * table_rows)
    )
    table = table_shape.table

    header_color = RGBColor(0x00, 0xB0, 0x50)
    row_color_light = RGBColor(0xE6, 0xF4, 0xEA)

    for row_idx in range(table_rows):
        for col_idx in range(table_cols):
            cell = table.cell(row_idx, col_idx)
            val = df.columns[col_idx] if row_idx == 0 else df.iloc[row_idx - 1, col_idx]
            if isinstance(val, float):
                val = round(val, 1)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(14 if row_idx == 0 else 12)
            para.font.name = 'Calibri'
            para.font.bold = row_idx == 0
            para.font.color.rgb = RGBColor(255, 255, 255) if row_idx == 0 else RGBColor(0, 0, 0)
            para.alignment = PP_ALIGN.CENTER if row_idx == 0 else PP_ALIGN.LEFT

            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = header_color if row_idx == 0 else (row_color_light if row_idx % 2 == 1 else RGBColor(255, 255, 255))

            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)

def add_bar_chart_slide(prs, grouped):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(8), Inches(0.5))
    title_box.text_frame.text = "SESSIONI ALL CHANNEL 2025"
    title_box.text_frame.paragraphs[0].font.size = Pt(20)
    title_box.text_frame.paragraphs[0].font.bold = True

    chart_data = CategoryChartData()
    chart_data.categories = list(grouped.index)
    for series_name in grouped.columns:
        values = [round(val, 1) for val in grouped[series_name].tolist()]
        chart_data.add_series(series_name, values)

    x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(5)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(12)
    data_labels.font.color.rgb = RGBColor(0, 0, 0)
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

# --- Crea presentazione ---
prs = Presentation()

# Slide 1: titolo
add_title_slide(prs, "Lievitosohn", "Report Febbraio 2025", black_bg=True)

# Slide 2: overview
add_title_slide(prs, "OVERVIEW")

# Slide 3: vuota
add_blank_slide(prs)

# Slide 4: bar chart
add_bar_chart_slide(prs, grouped)

# Slide 5: vuota
add_blank_slide(prs)

# Slide 6: tabella sessioni
slide_sessions = prs.slides.add_slide(prs.slide_layouts[5])
df_limited = prepare_table(df_sessions.head(20))
add_formatted_table(slide_sessions, df_limited)

# Slide 7: titolo MEDIA ORGANIC
add_title_slide(prs, "MEDIA ORGANIC")

# Slide 8-11: micro-moments
tables = [df_brand, df_product, df_need, df_best]

for df in tables:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    table_data = prepare_table(df)
    add_formatted_table(slide, table_data)

# Slide 12: titolo MEDIA PAID
add_title_slide(prs, "MEDIA PAID")

# Slide 13: vuota
add_blank_slide(prs)

# Slide 14: vuota con sfondo nero
add_blank_slide(prs, black_bg=True)

# --- Salvataggio ---
output_path = "seo_lievitosohn_report_feb25_starter.pptx"
prs.save(output_path)
print("✅ Presentazione salvata con successo!")
