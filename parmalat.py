from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION, XL_TICK_MARK
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd

template = "/data/template"
font_standard = "Gotham HTF"
font_light = "Gotham HTF Light" 
font_medium = "Gotham HTF Medium"

prs = Presentation(template)

content2023 = pd.read_csv("/data/parmalat_content2023.csv")
content2024 = pd.read_csv("/data/parmalat_content2024.csv")
trafficoOrganico2425 = pd.read_csv("/data/parmalat_traffico_organico_clic2425.csv")
trafficoOrganico2425 = trafficoOrganico2425.fillna("")
tables_data = {
    'parmalat': '/data/parmalat_focus_miglioramento_parmalat.csv',
    'chef': '/data/parmalat_focus_miglioramento_chef.csv',
    'zymil': '/data/parmalat_focus_miglioramento_zymil.csv',
    'santal': '/data/parmalat_focus_miglioramento_santal.csv',
    'cluster1': '/data/parmalat_focus_cluster1.csv',
    'cluster2': '/data/parmalat_focus_cluster2.csv'
}

mesi = trafficoOrganico2425["Mese"].tolist()
y2024 = trafficoOrganico2425["2024"].tolist()
y2025 = trafficoOrganico2425["2025"].tolist()
chart_data = CategoryChartData()
chart_data.categories = mesi
chart_data.add_series("2024", y2024)
chart_data.add_series("2025", y2025)

def add_clic_chart(data, slide):
    x, y, cx, cy = Cm(3.36), Cm(9.36), Cm(60.98), Cm(24.15)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, data
    )
    chart = graphic_frame.chart

    # Colori personalizzati
    color_map = {
        "2024": RGBColor(217, 217, 217), 
        "2025": RGBColor(64, 212, 48)
    }

    for series in chart.series:
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = color_map.get(series.name, RGBColor(0, 0, 0))

        # Etichette sopra ogni barra
        series.has_data_labels = True
        labels = series.data_labels
        labels.show_value = True
        labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
        txPr = labels._element.get_or_add_txPr()
        txPr.bodyPr.set('rot', '-5400000')  # Rotazione di -90 gradi
        labels.font.size = Pt(18)
        labels.font.color.rgb = RGBColor(64, 64, 64)

    # Asse Y
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.maximum_scale = 25000
    chart.value_axis.minimum_scale = 0
    chart.value_axis.tick_labels.font.size = Pt(20)

    # Asse X
    chart.category_axis.tick_labels.font.size = Pt(20)

    # Legenda
    chart.has_legend = True
    chart.legend.font.size = Pt(24)
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.font.name = font_light
    chart.legend.include_in_layout = False
    chart.value_axis.maximum_scale = 700000
    chart.value_axis.minimum_scale = 0

    # Griglia principale asse Y
    gridlines = chart.value_axis.major_gridlines
    line = gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(200, 200, 200)  # Grigio chiaro

    axis_line = chart.value_axis.format.line
    axis_line.fill.background()

    category_axis = chart.category_axis

    # 1. Disabilita i tick
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.minor_tick_mark = XL_TICK_MARK.NONE

    # 2. Colora la linea dell'asse
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = RGBColor(200, 200, 200)  # stesso colore griglia


def create_2023_table(slide, df):
    rows, cols = df.shape
    x, y, cx, cy = Cm(7.96), Cm(6.71), Cm(52.58), Cm(25.96)  # Posizione e dimensioni tabella

    table_shape = slide.shapes.add_table(rows + 1, cols, x, y, cx, cy)
    table = table_shape.table

    # Dimensioni riga header
    header_height = round(Cm(3.39))
    body_total_height = round(cy) - header_height
    body_row_height = round(body_total_height / rows)

    # Imposta header
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        if col_idx < 2:
            cell.text = ""
        else:
            cell.text = col_name.upper()
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Grigio header

        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(22)
        para.font.name = font_light
        para.alignment = PP_ALIGN.CENTER
        para.font.color.rgb = RGBColor(0, 0, 0)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE


    # Corpo tabella
    for row_idx, row in df.iterrows():
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)

            # Evita di scrivere 'nan' o 'NaN'
            if pd.isna(value):
                value = ""  # Lascia vuoto

            cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(32)
            para.font.name = font_light
            para.font.bold = False
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if col_idx < 2:
                para.alignment = PP_ALIGN.LEFT
            else:
                para.alignment = PP_ALIGN.CENTER

            # Sfondo bianco standard
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    


    # Merge celle per Brand
    table.cell(1, 0).merge(table.cell(2, 0))  # CHEF
    table.cell(3, 0).merge(table.cell(5, 0))  # ZYMIL
    table.cell(6, 0).merge(table.cell(8, 0))  # SANTAL

    # Colorazioni blocchi Brand
    chef_color = RGBColor(255, 230, 153)     # Giallo chiaro
    zymil_color = RGBColor(217, 217, 242)    # Azzurro chiaro
    santal_color = RGBColor(198, 224, 180)   # Verde chiaro
    total_color = RGBColor(64, 212, 48)       # Verde brillante

    # Colora celle Brand + Attività
    for idx in [1, 2]:
        for col in range(2):
            cell = table.cell(idx, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = chef_color

    for idx in [3, 4, 5]:
        for col in range(2):
            cell = table.cell(idx, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = zymil_color

    for idx in [6, 7, 8]:
        for col in range(2):
            cell = table.cell(idx, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = santal_color

    # Colora riga Totale
    for col_idx in range(cols):
        cell = table.cell(9, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = total_color
        para = cell.text_frame.paragraphs[0]
        para.font.bold = False
        para.font.color.rgb = RGBColor(0, 0, 0)
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # white background first two col of last row
    for col_idx in range(2):
        cell = table.cell(rows, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

def create_2024_table(slide, df):
    rows, cols = df.shape
    x, y, cx, cy = Cm(7.96), Cm(6.71), Cm(52.58), Cm(25.96)  # Posizione e dimensioni tabella

    table_shape = slide.shapes.add_table(rows + 1, cols, x, y, cx, cy)
    table = table_shape.table

    # Dimensioni riga header
    header_height = round(Cm(3.39))
    body_total_height = round(cy) - header_height
    body_row_height = round(body_total_height / rows)

    # Imposta header
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        if col_idx < 2:
            cell.text = ""
        else:
            cell.text = col_name.upper()
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(22)
        para.font.name = font_light
        para.alignment = PP_ALIGN.CENTER
        para.font.color.rgb = RGBColor(0, 0, 0)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Corpo tabella
    for row_idx, row in df.iterrows():
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)

            if pd.isna(value):
                value = ""  # Lascia vuoto

            cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(32)
            para.font.name = font_light
            para.font.bold = False
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if col_idx < 2:
                para.alignment = PP_ALIGN.LEFT
            else:
                para.alignment = PP_ALIGN.CENTER

            # Sfondo bianco standard
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Definizione colori
    chef_color = RGBColor(255, 230, 153)
    zymil_color = RGBColor(217, 217, 242)
    santal_color = RGBColor(198, 224, 180)
    total_color = RGBColor(64, 212, 48)

    # Applica colorazione righe per Brand
    brand_colors = {
        "CHEF": chef_color,
        "ZYMIL": zymil_color,
        "SANTAL": santal_color
    }

    for row_idx, row in df.iterrows():
        brand = str(row["Brand"]).strip().upper()
        for col_idx in range(2):  # Prime due colonne
            cell = table.cell(row_idx + 1, col_idx)
            if brand in brand_colors:
                cell.fill.solid()
                cell.fill.fore_color.rgb = brand_colors[brand]

    # Merge celle nella prima colonna per brand
    table.cell(1, 0).merge(table.cell(2, 0))  # CHEF: righe 1 e 2
    table.cell(3, 0).merge(table.cell(4, 0))  # ZYMIL: righe 3 e 4
    table.cell(5, 0).merge(table.cell(6, 0))  # SANTAL: righe 5 e 6

    # Colora riga Totale
    total_row_idx = rows
    for col_idx in range(cols):
        cell = table.cell(total_row_idx, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = total_color
        para = cell.text_frame.paragraphs[0]
        para.font.bold = False
        para.font.color.rgb = RGBColor(0, 0, 0)
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Sfondo bianco prime 2 colonne della riga Totale
    for col_idx in range(2):
        cell = table.cell(total_row_idx, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

def clean_header(text):
    """Pulisce testo header: rimuove \n, spazi extra, trim."""
    return ' '.join(text.replace('\n', ' ').split()).strip()

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def clean_header(text):
    """Pulisce testo header: rimuove \n, spazi extra, trim."""
    return ' '.join(text.replace('\n', ' ').split()).strip()

def fill_table_from_df(slide, df, header_row_idx=0):
    """
    Riempi una tabella esistente nella slide usando i dati di un DataFrame,
    formattando testo, allineamento e colori.
    
    Args:
        slide: oggetto slide di python-pptx.
        df: pandas DataFrame con i dati da inserire.
        header_row_idx: indice della riga che contiene l'header nella tabella (default 0).
    """
    # Trova la prima tabella nella slide
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break
    
    if table is None:
        raise ValueError("Nessuna tabella trovata nella slide.")

    # Leggi i nomi colonne della tabella e puliscili
    header = []
    for col_idx in range(len(table.columns)):
        cell = table.cell(header_row_idx, col_idx)
        header_text = clean_header(cell.text)
        header.append(header_text)

    # Crea la mappa: indice colonna tabella -> nome colonna CSV
    df_columns_clean = [clean_header(c) for c in df.columns]
    
    col_map = {}
    for col_idx, header_name in enumerate(header):
        if header_name in df_columns_clean:
            matched_col = df.columns[df_columns_clean.index(header_name)]
            col_map[col_idx] = matched_col
        else:
            print(f"[Warning] Colonna '{header_name}' non trovata nel DataFrame!")

    cols = len(table.columns)

    # Riempie le celle della tabella (dopo l'header)
    for row_idx in range(1, len(table.rows)):
        df_row_idx = row_idx - 1
        if df_row_idx >= len(df):
            break
        for col_idx, csv_col_name in col_map.items():
            value = df.iloc[df_row_idx][csv_col_name]
            cell = table.cell(row_idx, col_idx)
            if pd.isna(value):
                cell.text = ""
            else:
                cell.text = str(value)

            # Formattazione testo
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(24)
            para.font.name = font_light  # Puoi cambiare se vuoi un font diverso

            # Allineamento
            if col_idx == 0:
                para.alignment = PP_ALIGN.LEFT
            else:
                para.alignment = PP_ALIGN.CENTER

            # Colore testo per ultima colonna (Δ)
            if col_idx == cols - 1:  # ultima colonna
                value_str = str(value)
                if "+" in value_str:
                    para.font.color.rgb = RGBColor(0, 176, 80)  # Verde
                elif "-" in value_str:
                    para.font.color.rgb = RGBColor(255, 0, 0)    # Rosso
                else:
                    para.font.color.rgb = RGBColor(0, 0, 0)      # Nero
            else:
                para.font.color.rgb = RGBColor(0, 0, 0)  # Nero standard per tutte le altre colonne


from pptx.util import Cm

def insert_image(slide, image_path, left_cm=2, top_cm=2, width_cm=None, height_cm=None):
    """
    Inserisce un'immagine in una slide PowerPoint.
    
    Args:
        slide: oggetto slide di python-pptx.
        image_path: percorso all'immagine da inserire.
        left_cm: distanza dal bordo sinistro (in centimetri).
        top_cm: distanza dal bordo superiore (in centimetri).
        width_cm: larghezza desiderata (in centimetri, opzionale).
        height_cm: altezza desiderata (in centimetri, opzionale).
    """
    left = Cm(left_cm)
    top = Cm(top_cm)
    
    # Se specificati width o height
    if width_cm and height_cm:
        width = Cm(width_cm)
        height = Cm(height_cm)
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    else:
        # Inserisce con dimensioni originali
        slide.shapes.add_picture(image_path, left, top)



img_paths = {
    'focus1': 'focus1.png',
    'focus2': 'focus2.png',
    'focus3': 'focus3.png'
}



for i, slide in enumerate(prs.slides):
    if i == 3:
        # Add chart to the current slide
        add_clic_chart(chart_data, slide)
    elif i == 5:
    # Create a new table on the current slide
        create_2023_table(slide, content2023)
    elif i == 6:
        # Create a new table on the current slide
        create_2024_table(slide, content2024)
    elif i in [7, 8, 9, 10, 11, 12]:
        # Create a new table on the current slide
        brand = list(tables_data.keys())[i - 7]
        df = pd.read_csv(tables_data[brand])
        df = df.fillna("")
        fill_table_from_df(slide, df, header_row_idx=0)
    elif i in [14, 15, 16]:
        # Add image to the current slide
        img_path = img_paths[f'focus{i-13}']
        insert_image(slide, img_path, left_cm=5, top_cm=12.48, width_cm=34.38, height_cm=19.42)
    

# To get shapes in your slides
slides = [slide for slide in prs.slides]
shapes = []
for slide in slides:
    for shape in slide.shapes:
        shapes.append(shape)

def replace_text(replacements: dict, shapes: list):
    """Takes dict of {match: replacement, ... } and replaces all matches.
    Currently not implemented for charts or graphics.
    """
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            cur_text = run.text
                            new_text = cur_text.replace(str(match), str(replacement))
                            run.text = new_text
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if match in cell.text:
                            new_text = cell.text.replace(match, replacement)
                            cell.text = new_text

replace_text({
    '{{MESE}}': 'MARZO',
    '{{ANNO2PREC}}': '2023',
    '{{ANNOPREC}}': '2024',
    '{{ANNO}}': '2025',
    '{{ANNI}}': '2024-2025',
    '{{FOCUS1}}': 'ALIMENTI RICCHI DI POTASSIO',
    '{{FOCUS2}}': 'PROTEINE DEL LATTE',
    '{{FOCUS3}}': 'ANANAS BENEFICI',
    }, shapes) 


prs.save("parmalat_report.pptx")
