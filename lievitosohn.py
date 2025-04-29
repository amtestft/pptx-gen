from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION, XL_TICK_MARK
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd

template = "/data/template"

font_standard = "Gotham HTF"
font_light = "Gotham HTF Light" 
font_medium = "Gotham HTF Medium"

obiettivo_2024 = "149.700"

trafficAllSources = pd.read_csv("/data/traffic_all_sources.csv")
mesi = trafficAllSources["Mese"].tolist()
organic = trafficAllSources["Organic"].tolist()
other = trafficAllSources["Other"].tolist()
paid = trafficAllSources["Paid"].tolist()
totali = trafficAllSources["Total"].tolist()
chart_data = CategoryChartData()
chart_data.categories = mesi
chart_data.add_series("Organic", organic)
chart_data.add_series("Other", other)
chart_data.add_series("Paid", paid)

# Prepara dati per il grafico
clic_df = pd.read_csv("/data/clic_organici.csv")
clic_data = CategoryChartData()
clic_data.categories = clic_df["Mese"]
clic_data.add_series("2023", clic_df["2023"])
clic_data.add_series("2024", clic_df["2024"])
clic_data.add_series("Obiettivo di Progetto", clic_df["Obiettivo di Progetto"])

def edit_slide_2(slide, title):
    title_box = slide.shapes.add_textbox(Cm(3.37), Cm(10.2), Cm(51.13), Cm(12.25)) # left, top, width, height
    text_frame = title_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title
    p.font.size = Pt(124)
    p.font.name = font_medium
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def edit_slide_3and5(slide, title):
    title_box = slide.shapes.add_textbox(Cm(0), Cm(16.62), Cm(67.73), Cm(4.83)) # left, top, width, height
    text_frame = title_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title
    p.font.size = Pt(120)
    p.font.name = font_medium
    p.font.bold = True
    p.font.color.rgb = RGBColor(0,0,0)
    p.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_traffic_chart(chart_data, slide):
    # Aggiungi il grafico
    x, y, cx, cy = Cm(3.61), Cm(5.7), Cm(34.75), Cm(29.54) # left, top, height, width
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data
    )

    # 6. Imposta i colori: verde, blu, giallo
    chart = graphic_frame.chart
    custom_colors = {
        "Organic": RGBColor(0, 153, 102),  # Verde
        "Other": RGBColor(0, 51, 153),     # Blu
        "Paid": RGBColor(255, 204, 0),     # Giallo
    }

    for series in chart.series:
        series_name = series.name
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = custom_colors.get(series_name, RGBColor(0, 0, 0))  # Fallback nero se non trovato
        series.format

    # legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.maximum_scale = 110000
    chart.value_axis.minimum_scale = 0

    plot = chart.plots[0]
    plot.gap_width = 10

    # Griglia principale asse Y
    gridlines = chart.value_axis.major_gridlines
    line = gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(200, 200, 200)  # Grigio chiaro

    axis_line = chart.value_axis.format.line
    axis_line.fill.background()

    # Font asse Y
    chart.value_axis.tick_labels.font.size = Pt(18)
    chart.value_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)

    # Font asse X
    chart.category_axis.tick_labels.font.size = Pt(14)
    chart.category_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)

    # Font legenda
    chart.has_legend = True
    chart.legend.font.size = Pt(24)
    chart.legend.font.color.rgb = RGBColor(0, 0, 0)

    # font tutto il grafico
    chart.font.name = font_light

    category_axis = chart.category_axis

    # 1. Disabilita i tick
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.minor_tick_mark = XL_TICK_MARK.NONE

    # 2. Colora la linea dell'asse
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = RGBColor(200, 200, 200)  # stesso colore griglia

def prepare_pie_and_table_data_from_df(df):
    """
    df deve essere trafficAllSources: deve contenere le colonne
    ['Mese', 'Organic', 'Other', 'Paid', 'Total']
    """
    # Calcolo totali per il grafico a torta
    total_paid = df["Paid"].sum()
    total_organic = df["Organic"].sum()
    total_other = df["Other"].sum()

    pie_chart_data = CategoryChartData()
    pie_chart_data.categories = ["Paid", "Organic", "Other"]
    pie_chart_data.add_series('Traffico 2024', (total_paid, total_organic, total_other))

    # Preparo la tabella dettagliata
    table_data = []
    for idx, row in df.iterrows():
        mese = row["Mese"].capitalize()
        paid = row["Paid"]
        organic = row["Organic"]
        other = row["Other"]
        totale = paid + organic + other

        perc_paid = f"{round(100 * paid / totale)}%"
        perc_organic = f"{round(100 * organic / totale)}%"
        perc_other = f"{round(100 * other / totale)}%"

        table_data.append([mese, totale, perc_paid, perc_organic, perc_other])
    
    return pie_chart_data, table_data

pie_data, table_data = prepare_pie_and_table_data_from_df(trafficAllSources)

def add_pie_chart(pie_data, slide):
    from pptx.util import Cm, Pt
    from pptx.enum.chart import XL_LEGEND_POSITION

    x, y, cx, cy = Cm(42.73), Cm(5.7), Cm(19.08), Cm(11.77)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE_EXPLODED, x, y, cx, cy, pie_data
    )
    
    chart = graphic_frame.chart

    # Imposta colori personalizzati
    custom_colors_list = [
        RGBColor(255, 204, 0),     # Paid → Giallo
        RGBColor(0, 153, 102),     # Organic → Verde
        RGBColor(0, 51, 153),      # Other → Blu
    ]

    for idx, point in enumerate(chart.series[0].points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = custom_colors_list[idx]  # Basato solo su ordine!
        
    ''' # Esplodi spicchi piccoli
        if label != "Paid":
            point.explode = True'''

    # Aggiungi etichette
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.show_percentage = True
    data_labels.show_category_name = False
    data_labels.show_value = False
    data_labels.font.size = Pt(20)
    data_labels.font.name = font_light

    # Imposta legenda
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(24)
    chart.legend.font.name = font_standard


def add_traffic_table(table_data, slide):
    from pptx.util import Cm, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

    import pandas as pd
    df = pd.DataFrame(table_data, columns=["Mese", "Totale Visite", "% Paid", "% Organic", "% Other"])

    rows, cols = df.shape
    table_shape = slide.shapes.add_table(rows + 1, cols, Cm(40.41), Cm(19.71), Cm(23.71), Cm(14))  # left, top, width, height
    table = table_shape.table

    # Altezza header
    header_height = round(Cm(2))  # proporzionale al tuo esempio
    body_total_height = round(Cm(16)) - header_height
    body_row_height = round(body_total_height / rows)

    # Header
    table.rows[0].height = header_height

    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(21)
        para.font.name = "Arial"  # oppure font_medium
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(64, 212, 48)  # Verde header

    # Corpo tabella
    for row_idx, row in df.iterrows():
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            para.font.name = "Arial"  # oppure font_light
            para.font.size = Pt(20)
            para.font.bold = False

            # Sfondo bianco
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

            # Allineamento
            if col_idx == 0:
                para.alignment = PP_ALIGN.LEFT
            else:
                para.alignment = PP_ALIGN.CENTER

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            para.font.color.rgb = RGBColor(0, 0, 0)



def edit_slide_4(slide, title, chart_data, pie_data, table_data):
    # titolo
    title_box = slide.shapes.add_textbox(Cm(2.69), Cm(1.8), Cm(62.35), Cm(2.49)) # left, top, width, height
    text_frame = title_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title
    p.font.size = Pt(44)
    p.font.name = font_standard
    p.font.bold = False
    p.font.color.rgb = RGBColor(0,0,0)
    text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

    # add chart
    add_traffic_chart(chart_data, slide)
    add_pie_chart(pie_data, slide)
    add_traffic_table(table_data, slide)


def add_clic_chart(data, slide):
    x, y, cx, cy = Cm(4.1), Cm(8.26), Cm(57.88), Cm(26.59)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, data
    )
    chart = graphic_frame.chart

    # Colori personalizzati
    color_map = {
        "2023": RGBColor(217, 217, 217),
        "2024": RGBColor(64, 212, 48),
        "Obiettivo di Progetto": RGBColor(197, 224, 180),
    }

    for series in chart.series:
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = color_map.get(series.name, RGBColor(0, 0, 0))

        # Etichette sopra ogni barra
        series.has_data_labels = True
        labels = series.data_labels
        labels.show_value = True
        labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
        labels.font.size = Pt(18)
        txPr = labels._element.get_or_add_txPr()
        txPr.bodyPr.set('rot', '-5400000')  # Rotazione di -90 gradi
        labels.font.color.rgb = RGBColor(64, 64, 64)

    # Asse Y
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.maximum_scale = 25000
    chart.value_axis.minimum_scale = 0
    chart.value_axis.tick_labels.font.size = Pt(20)

    # Asse X
    chart.category_axis.tick_labels.font.size = Pt(20)

    # Legenda
    chart.has_legend = True
    chart.legend.font.size = Pt(24)
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.font.name = font_light
    chart.legend.include_in_layout = False
    chart.value_axis.maximum_scale = 25000
    chart.value_axis.minimum_scale = 0

    # Griglia principale asse Y
    gridlines = chart.value_axis.major_gridlines
    line = gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(200, 200, 200)  # Grigio chiaro

    axis_line = chart.value_axis.format.line
    axis_line.fill.background()

    category_axis = chart.category_axis

    # 1. Disabilita i tick
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.minor_tick_mark = XL_TICK_MARK.NONE

    # 2. Colora la linea dell'asse
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = RGBColor(200, 200, 200)  # stesso colore griglia


def edit_slide_6(slide, title, clic_data):
    # titolo
    title_box = slide.shapes.add_textbox(Cm(2.69), Cm(1.8), Cm(62.35), Cm(2.49)) # left, top, width, height
    text_frame = title_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title
    p.font.size = Pt(44)
    p.font.name = font_standard
    p.font.bold = False
    p.font.color.rgb = RGBColor(0,0,0)
    text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

    # textbox Obiettivo 2024
    text_box = slide.shapes.add_textbox(Cm(55.63), Cm(3.86), Cm(8.66), Cm(1.71)) # left, top, width, height
    text_frame = text_box.text_frame
    text_frame.margin_top = Cm(0.25)
    text_frame.margin_bottom = Cm(0.25)
    p = text_frame.paragraphs[0]
    p.text = "Obiettivo 2024"
    p.font.size = Pt(28)
    p.font.name = font_medium
    p.font.bold = False
    p.font.color.rgb = RGBColor(0,0,0)
    p.alignment = PP_ALIGN.CENTER

    # textbox numero obiettivo 2024
    text_box = slide.shapes.add_textbox(Cm(55.69), Cm(4.87), Cm(8.66), Cm(2.57)) # left, top, width, height
    text_frame = text_box.text_frame
    text_frame.margin_top = Cm(0.25)
    text_frame.margin_bottom = Cm(0.25)
    p = text_frame.paragraphs[0]
    p.text = obiettivo_2024
    p.font.size = Pt(48)
    p.font.name = font_standard
    p.font.bold = False
    p.font.color.rgb = RGBColor(64,212,48)
    p.alignment = PP_ALIGN.CENTER

    # add chart
    add_clic_chart(clic_data, slide)

    '''# Inserisci tutti i loghi
    google_updates, posizioni_mesi, offset_serie, google_logo_path, google_logo_width, google_logo_height, google_logo_pos_y = get_google_updates_info()
    for logo in google_updates:
        mese = logo['mese']
        serie = logo['serie']

        # posizione centrale del mese + offset della serie
        pos_x = Cm(posizioni_mesi[mese]) + offset_serie[serie]

        slide.shapes.add_picture(google_logo_path, pos_x, google_logo_pos_y, width=google_logo_width, height=google_logo_height)


    # 1. Disegna la linea verticale
    line = slide.shapes.add_shape(
        MSO_CONNECTOR.STRAIGHT,
        left=Cm(32.91),
        top=Cm(8.4),
        width=Pt(0),    # larghezza 0 --> linea verticale
        height=Cm(20.2)
    )

    # Setta colore e stile della linea
    line.line.color.rgb = RGBColor(191, 191, 191)  # Nero
    line.line.width = Pt(1.5)                  # Spessore linea

    textbox = slide.shapes.add_textbox(
        left=Cm(28.65),
        top=Cm(6.25),
        width=Cm(8.52),
        height=Cm(2.39)
    )
    text_frame = textbox.text_frame
    text_frame.clear()

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "GO LIVE 5/06\nRework contenuti"  # <-- \n qui!
    run.font.size = Pt(22)
    run.font.name = font_medium
    run.font.color.rgb = RGBColor(191, 191, 191)  # tutto nero (se vuoi colori diversi vedi sotto)
    p.alignment = PP_ALIGN.CENTER

    # Bordo del rettangolo (opzionale)
    textbox.line.color.rgb = RGBColor(191, 191, 191)
    textbox.fill.background()  # Sfondo trasparente
    textbox.line.width = Pt(1.5)  # Spessore del bordo
    text_frame.margin_left, text_frame.margin_right = Cm(0.51), Cm(0.51)
    text_frame.margin_top, text_frame.margin_bottom = Cm(0.25), Cm(0.25)'''





def add_mm_slides(prs, slide_indices, slides_data):
    """
    Aggiorna le slide 7-10 con i dati dei Micro Moments.
    
    Args:
        prs: Presentation object
        slides_data: list of dictionaries, each with:
            - title: str
            - data_path: str
            - rework_flag: bool
            - rows_reworked: list of int
    """
    
    for idx, slide_index in enumerate(slide_indices):
        slide_info = slides_data[idx]
        slide = prs.slides[slide_index]

        # Pulisce forme non placeholder
        for shape in list(slide.shapes):
            if not shape.is_placeholder:
                sp = shape
                sp.element.getparent().remove(sp.element)

        # Aggiungi titolo
        title_box = slide.shapes.add_textbox(Cm(3.05), Cm(0.51), Cm(60.51), Cm(3.64))
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = slide_info["title"]
        p.font.size = Pt(44)
        p.font.bold = False
        p.font.name = font_standard
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.LEFT
        title_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

        # Carica dati CSV
        df = pd.read_csv(slide_info["data_path"])
        df = df.iloc[:, 1:]

        # Aggiungi tabella
        rows, cols = df.shape
        table = slide.shapes.add_table(rows+1, cols, Cm(3.05), Cm(7.34), Cm(60.51), Cm(24.83)).table # left, top, width, height

        header_height = round(Cm(4.57))
        body_total_height = round(Cm(24.83)) - header_height
        body_row_height = round(body_total_height / rows)


        # Setta header
        table.rows[0].height = header_height

        # Setta righe corpo
        for i in range(1, rows + 1):
            table.rows[i].height = body_row_height

        # Header
        for col_idx, col_name in enumerate(df.columns):
            cell = table.cell(0, col_idx)
            cell.text = col_name
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(24)
            para.font.font_name = font_medium
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(64, 212, 48)  # verde header
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Corpo tabella
        for row_idx, row in df.iterrows():
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx+1, col_idx)
                cell.text = str(value)
                para = cell.text_frame.paragraphs[0]
                para.font.name = font_light
                para.font.size = Pt(24)
                para.font.bold = False

                # Sfondo bianco
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

                # Allineamento
                if col_idx == 0:
                    para.alignment = PP_ALIGN.LEFT
                else:
                    para.alignment = PP_ALIGN.CENTER

                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                # Colore testo (solo ultima colonna)
                if col_idx == cols - 1:  # ultima colonna
                    value_str = str(value)
                    if "+" in value_str:
                        para.font.color.rgb = RGBColor(0, 176, 80)  # Verde
                    elif "-" in value_str:
                        para.font.color.rgb = RGBColor(255, 0, 0)    # Rosso
                    else:
                        para.font.color.rgb = RGBColor(0, 0, 0)      # Nero
                else:
                    para.font.color.rgb = RGBColor(0, 0, 0)          # Nero normale



        # Evidenziazione righe se necessario
        if slide_info["rework_flag"]:
            for row_idx in slide_info["rows_reworked"]:
                for col_idx in range(cols):
                    cell = table.cell(row_idx+1, col_idx)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(197, 224, 180)  # verde chiaro

            # Aggiunge la legenda "Rework"
            legenda_box = slide.shapes.add_textbox(Cm(5.74), Cm(34.22), Cm(18.77), Cm(1.15))
            legenda_frame = legenda_box.text_frame
            legenda_frame.clear()
            p = legenda_frame.paragraphs[0]
            p.text = "Rework Content Strategy 2024"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.name = font_light
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.margin_left, p.margin_right = Cm(0.13), Cm(0.13)
            p.margin_top, p.margin_bottom = Cm(0.06), Cm(0.06)

            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Cm(3.05), Cm(34.28), Cm(2.69), Cm(1.03)
            )
            fill = rect.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(197, 224, 180)  # Verde chiaro come nella tabella
            rect.line.fill.background()  # Nessun bordo nero attorno


prs = Presentation(template)

slides = prs.slides
for i, slide in enumerate(slides):
    if i == 1:
        edit_slide_2(slide, "REPORT DICEMBRE")
    elif i == 2:
        edit_slide_3and5(slide, "OVERVIEW")
    elif i == 3:
        edit_slide_4(slide, "TRAFFIC ALL SOURCES 2024", chart_data, pie_data, table_data)
    elif i == 4:
        edit_slide_3and5(slide, "MEDIA ORGANIC")
    elif i == 5:
        edit_slide_6(slide, "CLIC ORGANICI E OBIETTIVO 2024", clic_data)

slide_idxs = [6, 7, 8, 9]
slides_data = [
    {
        "title": "TOP MICRO-MOMENTS - BRAND",
        "data_path": "/data/topmm_brand.csv",
        "rework_flag": False,
        "rows_reworked": []
    },
    {
        "title": "TOP MICRO-MOMENTS – PRODOTTO",
        "data_path": "/data/topmm_prodotto.csv",
        "rework_flag": False,
        "rows_reworked": []
    },
    {
        "title": "TOP MICRO-MOMENTS – STATI DI BISOGNO",
        "data_path": "/data/topmm_needs.csv",
        "rework_flag": True,
        "rows_reworked": [0, 2, 3]
    },
    {
        "title": "BEST VOLUME MICRO-MOMENTS",
        "data_path": "/data/best_vol_mm.csv",
        "rework_flag": True,
        "rows_reworked": [2, 4]
    }
]
add_mm_slides(prs, slide_idxs, slides_data)

prs.save("lievitosohn_report.pptx")
